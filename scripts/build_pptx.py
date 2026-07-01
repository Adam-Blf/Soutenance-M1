#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generateur PPTX soutenance Bloc 1 & 2 (RNCP40875).
Charte visuelle EFREI.fr : #163767 / #FF43B8 / #051832 / #0C78B4.
Police Gilroy (natif Windows) + fallback Calibri.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

try:
    from PIL import Image
    HAVE_PIL = True
except Exception:
    HAVE_PIL = False

ROOT    = Path(__file__).resolve().parent.parent
ASSETS  = ROOT / "assets"
SCREENS = ASSETS / "screens"
OUT     = ROOT / "Soutenance_Bloc1-2_Adam_Emilien.pptx"

# ---- Charte EFREI.fr exacte ----
EFREI_NAVY   = RGBColor(0x16, 0x37, 0x67)   # #163767  primaire
EFREI_PINK   = RGBColor(0xFF, 0x43, 0xB8)   # #FF43B8  accent signature
EFREI_DARK   = RGBColor(0x05, 0x18, 0x32)   # #051832  fond sombre
EFREI_BLUE   = RGBColor(0x0C, 0x78, 0xB4)   # #0C78B4  secondaire
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xF8, 0xF9, 0xFF)
RULE         = RGBColor(0xD8, 0xE0, 0xF0)
TEXT_DARK    = RGBColor(0x0D, 0x1B, 0x33)
TEXT_MUTED   = RGBColor(0x55, 0x6B, 0x8A)
SUCCESS      = RGBColor(0x0C, 0xB4, 0x87)

# Couleurs par section (projet)
P1   = EFREI_NAVY    # Urban Data Explorer
P2   = EFREI_BLUE    # Maintenance Predictive
P3   = EFREI_PINK    # IA Pero

# Alias backward-compat
BLEU      = EFREI_NAVY
TEAL      = EFREI_BLUE
ROUGE     = EFREI_PINK
VERT      = SUCCESS
TEXTE     = TEXT_DARK
GRIS      = TEXT_MUTED
GRIS_CLR  = OFF_WHITE
GRIS_BORD = RULE
BLANC     = WHITE

# Police : Gilroy si installee, Calibri sinon
FONT_H = "Gilroy"     # titres & kickers
FONT_B = "Calibri"    # corps
FONT   = FONT_B

REPOS = {
    "urban":  "https://github.com/Adam-Blf/urban-data-explorer",
    "maint":  "https://github.com/Adam-Blf/maintenance-predictive-industrielle",
    "iapero": "https://github.com/Adam-Blf/ia-pero",
}
DEMO_URBAN = "https://urban.beloucif.com"

SW, SH   = 13.333, 7.5
EMU_IN   = 914400

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def _set_font(run, size, color, bold=False, italic=False, font=FONT_B):
    run.font.size      = Pt(size)
    run.font.color.rgb = color
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.name      = font


def textbox(s, text, l, t, w, h, size=18, color=TEXT_DARK, bold=False, italic=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_B, line_spacing=1.0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment    = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    _set_font(r, size, color, bold, italic, font)
    return tb


def bullets(s, items, l, t, w, h, size=15, color=TEXT_DARK, gap=7, bullet_color=EFREI_PINK):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after  = Pt(gap)
        p.line_spacing = 1.08
        rb = p.add_run()
        rb.text = "  "
        _set_font(rb, size - 3, bullet_color, bold=True)
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = lead; _set_font(r1, size, color, bold=True, font=FONT_H)
            r2 = p.add_run(); r2.text = rest; _set_font(r2, size, color)
        else:
            r = p.add_run(); r.text = it; _set_font(r, size, color)
    return tb


def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE, round_=None):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if round_ is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = round_
        except Exception:
            pass
    return sp


def shape_text(sp, text, size=13, color=WHITE, bold=True,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(6)
    tf.margin_top  = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    _set_font(r, size, color, bold, font=FONT_H)
    return sp


def header(s, kicker, title, color=EFREI_NAVY):
    """Header EFREI : bande coloree fine + kicker rose + titre navy."""
    bg(s, WHITE)
    rect(s, 0, 0, SW, 0.12, fill=color)
    textbox(s, kicker.upper(), 0.52, 0.2, 11.0, 0.3,
            size=8.5, color=EFREI_PINK, bold=True, font=FONT_H)
    textbox(s, title, 0.52, 0.46, 11.5, 0.65,
            size=24, color=EFREI_DARK, bold=True, font=FONT_H)
    rect(s, 0.52, 1.18, SW - 1.05, 0.018, fill=RULE)


def footer(s, idx, dark=False):
    c = RGBColor(0x88, 0x9A, 0xBB) if dark else TEXT_MUTED
    rect(s, 0.52, 7.1, SW - 1.05, 0.012, fill=RULE)
    textbox(s, "Adam Beloucif  &  Emilien Morice  |  RNCP 40875  |  EFREI Villejuif",
            0.52, 7.15, 9.0, 0.28, size=8.5, color=c, font=FONT_B)
    textbox(s, f"{idx} / 20", 12.2, 7.15, 0.8, 0.28,
            size=8.5, color=c, align=PP_ALIGN.RIGHT)


def badge(s, code, label, l, t, fill=EFREI_PINK):
    w = 0.52 + 0.1 * len(label)
    sp = rect(s, l, t, min(w, 4.0), 0.36,
              fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.5)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(7)
    tf.margin_top  = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run(); r1.text = code + "  "; _set_font(r1, 10.5, WHITE, bold=True, font=FONT_H)
    r2 = p.add_run(); r2.text = label;        _set_font(r2, 9.5, WHITE, font=FONT_B)
    return sp


def badges_row(s, pairs, l, t, fill=EFREI_PINK):
    x = l
    for code, label in pairs:
        sp = badge(s, code, label, x, t, fill=fill)
        x += sp.width / EMU_IN + 0.15


def button(s, text, url, l, t, w=2.7, h=0.5, fill=EFREI_NAVY, fg=WHITE):
    sp = rect(s, l, t, w, h, fill=fill,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.3)
    shape_text(sp, text, size=11.5, color=fg, bold=True)
    sp.click_action.hyperlink.address = url
    return sp


def image_fit(s, path, l, t, max_w, max_h, border=False):
    path = Path(path)
    if not path.exists():
        rect(s, l, t, max_w, max_h, fill=OFF_WHITE, line=RULE)
        textbox(s, "[image]", l, t + max_h / 2 - 0.2, max_w, 0.4,
                size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        return
    if HAVE_PIL:
        with Image.open(path) as im:
            iw, ih = im.size
        ar = iw / ih; box_ar = max_w / max_h
        if ar > box_ar: w = max_w;  h = max_w / ar
        else:           h = max_h;  w = max_h * ar
        l2 = l + (max_w - w) / 2
        t2 = t + (max_h - h) / 2
    else:
        w, h, l2, t2 = max_w, max_h, l, t
    if border:
        rect(s, l2 - 0.03, t2 - 0.03, w + 0.06, h + 0.06,
             fill=WHITE, line=RULE, line_w=0.8)
    s.shapes.add_picture(str(path), Inches(l2), Inches(t2), Inches(w), Inches(h))


def caption(s, text, l, t, w):
    textbox(s, text, l, t, w, 0.28, size=9.5, color=TEXT_MUTED,
            italic=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1 - Titre (fond EFREI_DARK, typographie Gilroy)
# ============================================================
s = slide(); bg(s, EFREI_DARK)

# Bande rose gauche signature EFREI
rect(s, 0, 0, 0.22, SH, fill=EFREI_PINK)

# Bande rose fine en haut
rect(s, 0, 0, SW, 0.08, fill=EFREI_PINK)

# Kicker
textbox(s, "SOUTENANCE DE FIN D'ETUDES  |  M1 DATA ENGINEERING & IA",
        0.55, 1.2, 12.0, 0.42, size=11.5,
        color=RGBColor(0x6B, 0xA3, 0xD8), bold=True, font=FONT_H)

# Titre principal
textbox(s, "Bloc 1 & 2", 0.55, 1.7, 12.0, 0.75,
        size=52, color=WHITE, bold=True, font=FONT_H, line_spacing=0.9)
textbox(s, "RNCP 40875", 0.55, 2.5, 12.0, 0.62,
        size=28, color=EFREI_PINK, bold=True, font=FONT_H)

# Sous-titre
textbox(s, "Expert en Ingenierie de Donnees",
        0.55, 3.35, 12.0, 0.5, size=18,
        color=RGBColor(0x9B, 0xBD, 0xE0))

# Separateur
rect(s, 0.55, 4.02, 5.8, 0.04, fill=EFREI_PINK)

# Auteurs
textbox(s, "Adam Beloucif  &  Emilien Morice",
        0.55, 4.18, 12.0, 0.55, size=22, color=WHITE, bold=True, font=FONT_H)
textbox(s, "EFREI Villejuif  |  Session 2025-2026",
        0.55, 4.82, 12.0, 0.38, size=13,
        color=RGBColor(0x6B, 0xA3, 0xD8))

# ============================================================
# SLIDE 2 - Contexte
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Introduction", "Contexte & problematique")
textbox(s, "Fil rouge  -  transformer la donnee en decision",
        0.52, 1.42, 12, 0.46, size=18, color=EFREI_NAVY, bold=True, font=FONT_H)
textbox(s, "Trois projets complementaires couvrant la chaine de valeur data complete.",
        0.52, 1.92, 12, 0.36, size=13.5, color=TEXT_MUTED)
maillons = [
    ("1  Data Engineering",
     "Construire l'architecture\nqui collecte, stocke et\nexpose la donnee.", P1),
    ("2  Data Science",
     "Modeliser pour predire\net outiller la decision\nmetier.", P2),
    ("3  IA Generative",
     "Comprendre le langage\nnaturel et generer une\nreponse personnalisee.", P3),
]
x = 0.52
for titre, corps, col in maillons:
    rect(s, x, 2.5, 3.95, 2.85, fill=OFF_WHITE, line=RULE)
    rect(s, x, 2.5, 3.95, 0.1, fill=col)
    textbox(s, titre, x + 0.22, 2.75, 3.5, 0.46,
            size=15.5, color=col, bold=True, font=FONT_H)
    textbox(s, corps, x + 0.22, 3.32, 3.5, 1.7,
            size=13.5, color=TEXT_DARK, line_spacing=1.18)
    x += 4.18
textbox(s, "Une donnee non exploitee n'a aucune valeur. Chaque projet repond a un besoin "
           "metier concret, avec preuves techniques verifiables dans le code.",
        0.52, 5.65, 12.2, 0.78, size=13, color=TEXT_DARK, italic=True)
footer(s, 2)

# ============================================================
# SLIDE 3 - Vue d'ensemble
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Introduction", "Vue d'ensemble des 3 projets")
cards = [
    ("Urban Data Explorer", "Bloc 1  |  C1.1 - C2.4",
     "Plateforme data du logement parisien. 24 sources Open Data, architecture medaillon, "
     "PostgreSQL etoile, Cassandra, Kafka, API FastAPI securisee, dashboard MapLibre.",
     "urban", P1),
    ("Maintenance Predictive", "Bloc 2  |  C3.1 - C4.3",
     "Prediction de panne machine sous 24 h depuis capteurs IoT. 4 modeles (dont Deep Learning), "
     "mesure CO2 CodeCarbon, dashboard decisionnel Streamlit.",
     "maint", P2),
    ("L'IA Pero", "Bloc 2  |  C5.1 - C5.3",
     "Recommandation de cocktails par IA semantique. SBERT 384 dims, guardrail cosinus 0.35, "
     "RAG + Gemini avec cache SQLite, interface Streamlit Speakeasy.",
     "iapero", P3),
]
y = 1.45
for titre, comp, desc, key, col in cards:
    rect(s, 0.52, y, 12.3, 1.55, fill=OFF_WHITE, line=RULE)
    rect(s, 0.52, y, 0.12, 1.55, fill=col)
    textbox(s, titre, 0.82, y + 0.1, 5.5, 0.44,
            size=17.5, color=col, bold=True, font=FONT_H)
    textbox(s, comp,  0.82, y + 0.56, 5.5, 0.3,
            size=10.5, color=TEXT_MUTED, bold=True)
    textbox(s, desc,  0.82, y + 0.88, 8.7, 0.58,
            size=12, color=TEXT_DARK, line_spacing=1.05)
    button(s, "Repo ->", REPOS[key], 10.1, y + 0.5, w=2.38, h=0.52, fill=col)
    y += 1.7
footer(s, 3)

# ============================================================
# SLIDE 4 - Contributions
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Introduction", "Contributions individuelles")
textbox(s, "Presentation collective, evaluation individualisee. Repartition par projet :",
        0.52, 1.42, 12, 0.36, size=13, color=TEXT_MUTED)
rows = [
    ("Projet", "Adam Beloucif", "Emilien Morice"),
    ("Urban Data Explorer",
     "Architecture data, API securisee, modele PG/Cassandra, pipeline",
     "Sourcing & qualite 24 sources, front cartographique, data prep"),
    ("Maintenance Predictive",
     "Modelisation 4 modeles, selection XGBoost, dashboard Streamlit",
     "EDA, preparation donnees, analyse des correlations"),
    ("L'IA Pero",
     "Backend RAG/SBERT, guardrail semantique, cache SQLite, evaluation",
     "Referentiel cocktails, scenarios d'usage, tests des seuils"),
]
col_x = [0.52, 3.42, 8.12]
col_w = [2.9, 4.7, 4.65]
y = 1.92; rh = [0.48, 1.02, 1.02, 1.02]
for ri, row in enumerate(rows):
    is_head = ri == 0
    for ci, cell in enumerate(row):
        fill = EFREI_DARK if is_head else (OFF_WHITE if ri % 2 else WHITE)
        sp = rect(s, col_x[ci], y, col_w[ci], rh[ri], fill=fill, line=RULE, line_w=0.75)
        tf = sp.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Pt(8)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = cell
        _set_font(r, 12.5 if is_head else 11.5,
                  WHITE if is_head else TEXT_DARK,
                  bold=(is_head or ci == 0), font=FONT_H if is_head else FONT_B)
    y += rh[ri]
textbox(s, "Chaque membre porte au moins un choix technique majeur et en connait les limites.",
        0.52, 6.52, 12, 0.45, size=11.5, color=TEXT_MUTED, italic=True)
footer(s, 4)

# ============================================================
# P1 - SLIDE 5 - Besoin & architecture
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 1  |  Urban Data Explorer", "Besoin & architecture", P1)
bullets(s, [
    ("Besoin metier  ",
     "explorer le marche du logement parisien en croisant 24 sources Open Data heterogenes."),
    ("Architecture medaillon  ",
     "Bronze (brut Parquet) -> Silver (normalise, geocode) -> Gold (datamarts)."),
    ("Stockage dual  ",
     "PostgreSQL en etoile pour l'analytique, Cassandra query-first pour le streaming."),
    ("Exposition  ",
     "API FastAPI securisee -> dashboard React / MapLibre."),
], 0.52, 1.45, 5.9, 4.2, size=14, gap=10, bullet_color=EFREI_PINK)
image_fit(s, SCREENS / "ude-light.png", 6.7, 1.5, 6.15, 4.35, border=True)
caption(s, "Dashboard Urban Data Explorer  -  choroplethe prix m2", 6.7, 5.9, 6.15)
button(s, "Demo live ->", DEMO_URBAN, 0.52, 6.4, w=2.4, h=0.5, fill=P1)
button(s, "Repo ->", REPOS["urban"], 3.1, 6.4, w=2.0, h=0.5, fill=EFREI_DARK)
footer(s, 5)

# ============================================================
# P1 - SLIDE 6 - Bases de donnees
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 1  |  Urban Data Explorer", "Bases de donnees  SQL vs NoSQL", P1)
badges_row(s, [("C1.1", "Relationnelle"), ("C1.2", "NoSQL")], 0.52, 1.3, fill=P1)
rect(s, 0.52, 1.95, 6.02, 4.22, fill=OFF_WHITE, line=RULE)
rect(s, 0.52, 1.95, 6.02, 0.5, fill=P1)
textbox(s, "PostgreSQL  -  schema en etoile", 0.72, 2.02, 5.6, 0.36,
        size=14.5, color=WHITE, bold=True, font=FONT_H)
bullets(s, [
    "Fait + dimensions (arrondissement, temps), FK et cles composees",
    "Index sur axes de requete -> p95 < 4 ms en jointure",
    "Contraintes NOT NULL / integrite referentielle",
    "Test de charge : scripts/test_load_postgres.py",
    "Besoin : agregats analytiques par arrondissement x mois",
], 0.76, 2.58, 5.5, 3.3, size=12.5, gap=7, bullet_color=EFREI_PINK)
rect(s, 6.82, 1.95, 5.98, 4.22, fill=OFF_WHITE, line=RULE)
rect(s, 6.82, 1.95, 5.98, 0.5, fill=P2)
textbox(s, "Cassandra  -  query-first", 7.02, 2.02, 5.6, 0.36,
        size=14.5, color=WHITE, bold=True, font=FONT_H)
bullets(s, [
    "Partition par event_type, clustering event_time DESC",
    "TTL 7 jours -> la donnee chaude expire seule",
    "Deux patterns d'acces : par type et par arrondissement",
    "Adapte aux evenements semi-structures a forte velocite",
    "PG aurait impose un schema rigide + purges manuelles",
], 7.06, 2.58, 5.5, 3.3, size=12.5, gap=7, bullet_color=P2)
footer(s, 6)

# ============================================================
# P1 - SLIDE 7 - Data Lake & streaming
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 1  |  Urban Data Explorer", "Data Lake & streaming temps reel", P1)
badges_row(s, [("C1.3", "Data Lake"), ("C2.2", "Streaming")], 0.52, 1.3, fill=P1)
bullets(s, [
    ("24 sources / 8 familles  ",
     "CSV, GeoJSON, API -> atterrissage Bronze en Parquet."),
    ("Silver  ",
     "normalisation des codes, geocodage point-in-polygon IRIS."),
    ("Gold  ",
     "datamarts dashboard & timeline, format colonnaire optimise."),
    ("Kafka  ",
     "producteur / consommateur temps reel evenement par evenement."),
    ("Micro-batch  ",
     "streaming/microbatch.py -> fenetres tumbling 10 s (count, moyenne)."),
], 0.52, 1.88, 6.0, 4.05, size=14, gap=9, bullet_color=EFREI_PINK)
image_fit(s, SCREENS / "infrastructures_stacked_bar.png", 6.88, 1.9, 5.9, 3.0, border=True)
image_fit(s, SCREENS / "geocoding_integrity.png",         6.88, 5.05, 5.9, 1.62, border=True)
caption(s, "Repartition des sources + integrite du geocodage", 6.88, 6.7, 5.9)
footer(s, 7)

# ============================================================
# P1 - SLIDE 8 - API & securite
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 1  |  Urban Data Explorer", "API & securite", P1)
badge(s, "C2.1", "API securisee", 0.52, 1.3, fill=P1)
sp_d1 = rect(s, 9.52, 1.26, 3.28, 0.48, fill=EFREI_PINK,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.3)
shape_text(sp_d1, "DEMO LIVE  -  90 s", size=12.5, color=WHITE)
bullets(s, [
    ("JWT HS256  ",
     "secret en variable d'env (jamais commite), roles viewer / admin."),
    ("401 vs 403  ",
     "authentification et autorisation differenciees."),
    ("Quotas par IP  ",
     "120 anonyme / 600 authentifie req/min -> 429 si depassement."),
    ("Swagger  ",
     "documentation interactive live sur /docs."),
    ("CORS  ",
     "restreint a l'origine du front, donnees Open Data (RGPD by design)."),
], 0.52, 1.95, 6.0, 3.62, size=13.5, gap=8, bullet_color=EFREI_PINK)
textbox(s, "Scenario demo", 0.52, 5.65, 6, 0.32, size=12.5, color=P1, bold=True, font=FONT_H)
textbox(s, "Dashboard -> clic 11e arr. -> comparaison 11e/16e -> /docs -> "
           "curl token -> 403 sans role admin.",
        0.52, 5.96, 6.2, 0.78, size=12, color=TEXT_DARK, italic=True)
image_fit(s, SCREENS / "ude-api.png", 6.88, 1.9, 5.9, 4.05, border=True)
button(s, "Repo ->", REPOS["urban"], 6.88, 6.18, w=2.55, h=0.5, fill=P1)
button(s, "Demo live ->", DEMO_URBAN, 9.62, 6.18, w=2.45, h=0.5, fill=EFREI_DARK)
footer(s, 8)

# ============================================================
# P1 - SLIDE 9 - Performance & limites
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 1  |  Urban Data Explorer", "Performance, resilience & limites", P1)
badges_row(s, [("C1.4", "Resilience"), ("C2.4", "Perf pipelines")], 0.52, 1.3, fill=P1)
bullets(s, [
    ("Metriques pipeline  ",
     "etl/metrics.py -> duree, lignes, octets, lignes/sec via GET /pipeline/metrics."),
    ("Optimisations  ",
     "Polars (Rust) vs pandas, Parquet colonnaire, lru_cache geocodage inverse."),
    ("Resilience testee  ",
     "scripts/test_resilience.py : kill Postgres -> fallback parquet -> restart."),
    ("Compose  ",
     "healthchecks, restart unless-stopped, depends_on service_healthy."),
], 0.52, 1.9, 12.25, 2.62, size=14, gap=9, bullet_color=EFREI_PINK)
rect(s, 0.52, 4.72, 12.28, 1.75,
     fill=RGBColor(0xFF, 0xF0, 0xF8), line=EFREI_PINK, line_w=1.0)
textbox(s, "Limites assumees & roadmap", 0.78, 4.88, 11, 0.38,
        size=14.5, color=EFREI_PINK, bold=True, font=FONT_H)
textbox(s, "Demo mono-noeud (replication_factor 1), pas de HA reelle. "
           "Chemin documente vers cluster 3 noeuds + replication Cassandra + API stateless replicable.",
        0.78, 5.3, 11.82, 1.0, size=13.5, color=TEXT_DARK, line_spacing=1.12)
footer(s, 9)

# ============================================================
# P2 - SLIDE 10 - Question metier
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 2  |  Maintenance Predictive", "Question metier & strategie IA", P2)
badge(s, "C4.1", "Strategie d'integration IA", 0.52, 1.3, fill=P2)
bullets(s, [
    ("Probleme  ",
     "une panne non planifiee coute 5 a 50 k EUR/h d'arret de ligne."),
    ("Corrective  ",
     "on repare apres la casse -> cout maximal, imprevu."),
    ("Preventive  ",
     "entretien calendaire -> sur-maintenance, pieces gaspillees."),
    ("Predictive (cible)  ",
     "predire la panne sous 24 h -> intervenir au bon moment."),
], 0.52, 1.92, 7.0, 3.0, size=14, gap=10, bullet_color=P2)
rect(s, 8.02, 1.92, 4.78, 3.82, fill=OFF_WHITE, line=RULE)
textbox(s, "Feuille de route", 8.26, 2.06, 4.3, 0.38,
        size=14.5, color=P2, bold=True, font=FONT_H)
bullets(s, [
    "Capteurs IoT existants -> collecte",
    "Modele predictif + seuil de decision",
    "Alerte responsable maintenance",
    "Boucle de reentrainement periodique",
], 8.26, 2.58, 4.4, 3.0, size=13, gap=10, bullet_color=P2)
textbox(s, "Cout corrective ~100 EUR/h  ->  predictive ~20 EUR/h",
        0.52, 5.38, 7, 0.4, size=14, color=P2, bold=True, font=FONT_H)
footer(s, 10)

# ============================================================
# P2 - SLIDE 11 - Donnees & EDA
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 2  |  Maintenance Predictive", "Donnees & analyse exploratoire", P2)
badges_row(s, [("C3.1", "Preparation"), ("C3.3", "EDA")], 0.52, 1.3, fill=P2)
bullets(s, [
    ("Dataset  ",
     "Kaggle industrial_machine_maintenance  -  24 042 lignes x 15 variables."),
    ("Cible  ",
     "failure_within_24h (classe rare -> desequilibre)."),
    ("Nettoyage  ",
     "imputation mediane (fit train), winsorisation outliers capteurs."),
    ("Insight cle  ",
     "vibration_rms et temperature_motor : variables les plus correlees a la panne."),
    ("Coherence physique  ",
     "l'usure mecanique genere vibration ET echauffement simultanement."),
    ("EDA  ",
     "6 etapes : distributions, boxplots outliers, heatmap correlations, desequilibre 88/12 %, scatter vibration x temperature, tendances temporelles."),
], 0.52, 1.9, 6.12, 4.05, size=14, gap=8, bullet_color=P2)
image_fit(s, SCREENS / "correlation_matrix.png", 6.88, 1.9, 5.9, 4.55, border=True)
caption(s, "Matrice de correlation des variables capteurs", 6.88, 6.5, 5.9)
footer(s, 11)

# ============================================================
# P2 - SLIDE 12 - Modeles & comparaison
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 2  |  Maintenance Predictive", "Modeles & evaluation comparative", P2)
badges_row(s, [("C4.2", "Modeles"), ("C4.3", "Eval + ecoresponsabilite")], 0.52, 1.3, fill=P2)
data = [
    ["Modele",             "F1",    "ROC-AUC", "Temps", "CO2"],
    ["LogReg (baseline)",  "0.865", "0.914",   "3.2 s", "0.3 mg"],
    ["Random Forest",      "0.910", "0.952",   "45 s",  "8.2 mg"],
    ["XGBoost (retenu)",   "0.928", "0.964",   "35 s",  "6.1 mg"],
    ["MLP 64-32-16 (DL)",  "0.896", "0.936",   "18 s",  "3.8 mg"],
]
cx = [0.52, 4.52, 6.12, 8.02, 9.62]
cw = [4.0, 1.6, 1.9, 1.6, 1.7]
ty = 1.96; rh = 0.6
for ri, row in enumerate(data):
    retained = ri == 3
    for ci, cell in enumerate(row):
        if ri == 0:       fill = EFREI_DARK
        elif retained:    fill = RGBColor(0xE0, 0xF0, 0xFF)
        else:             fill = OFF_WHITE if ri % 2 else WHITE
        sp = rect(s, cx[ci], ty, cw[ci], rh, fill=fill, line=RULE, line_w=0.75)
        tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Pt(7)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = cell
        clr = WHITE if ri == 0 else (P2 if retained else TEXT_DARK)
        _set_font(r, 12, clr, bold=(ri == 0 or retained), font=FONT_H if ri == 0 else FONT_B)
    ty += rh
bullets(s, [
    ("Selection  ",
     "score = F1 - 0.5 x sigma(F1 en cross-val) -> performance ET stabilite."),
    ("Anti-leakage  ",
     "pipeline sklearn fit sur train uniquement (ADR dedie)."),
    ("Ecoresponsabilite  ",
     "CodeCarbon : XGBoost 6.1 mg < RF 8.2 mg pour de meilleures perfs."),
], 0.52, 5.22, 12.25, 1.75, size=12.5, gap=6, bullet_color=P2)
footer(s, 12)

# ============================================================
# P2 - SLIDE 13 - Dashboard + DEMO
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 2  |  Maintenance Predictive", "Dashboard decisionnel", P2)
badge(s, "C3.2", "Dashboard interactif", 0.52, 1.3, fill=P2)
sp_d2 = rect(s, 9.52, 1.26, 3.28, 0.48, fill=EFREI_PINK,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.3)
shape_text(sp_d2, "DEMO LIVE  -  90 s", size=12.5, color=WHITE)
bullets(s, [
    ("Cible  ",      "responsable maintenance, pas le data scientist."),
    ("Simulation  ", "sliders capteurs -> prediction temps reel + probabilite."),
    ("Explication  ","top variables influentes en langage metier."),
    ("Distinct  ",   "des visuels EDA du rapport (consigne du sujet)."),
], 0.52, 1.92, 5.9, 3.22, size=13.5, gap=9, bullet_color=P2)
textbox(s, "Scenario : vibration haute + temperature haute -> risque eleve -> feature importance.",
        0.52, 5.28, 6.0, 0.78, size=12, color=TEXT_MUTED, italic=True)
image_fit(s, SCREENS / "maintenance-dashboard-live.png", 6.88, 1.9, 5.9, 4.05, border=True)
button(s, "Ouvrir le repo ->", REPOS["maint"], 6.88, 6.15, w=5.9, h=0.5, fill=P2)
footer(s, 13)

# ============================================================
# P2 - SLIDE 14 - Limites
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 2  |  Maintenance Predictive", "Limites & bonus realises", P2)
rect(s, 0.52, 1.52, 6.02, 4.68,
     fill=RGBColor(0xFF, 0xF0, 0xF8), line=EFREI_PINK)
textbox(s, "Limites", 0.76, 1.68, 5.5, 0.38,
        size=16.5, color=EFREI_PINK, bold=True, font=FONT_H)
bullets(s, [
    "Donnees simulees (pas de bruit capteur reel)",
    "Pas de dimension temporelle exploitee",
    "Un LSTM sur sequences serait la suite logique",
    "Derive en production -> monitoring + reentrainement",
], 0.76, 2.2, 5.5, 3.75, size=13.5, gap=12,
    bullet_color=EFREI_PINK)
rect(s, 6.82, 1.52, 5.98, 4.68,
     fill=RGBColor(0xE8, 0xF4, 0xFF), line=P2)
textbox(s, "Bonus deja realises", 7.06, 1.68, 5.5, 0.38,
        size=16.5, color=P2, bold=True, font=FONT_H)
bullets(s, [
    "Multiclass : type de panne predit",
    "Regression : RUL (duree de vie restante)",
    "Calibration des probabilites",
    "Tuning des hyperparametres justifie",
], 7.06, 2.2, 5.5, 3.75, size=13.5, gap=12, bullet_color=P2)
footer(s, 14)

# ============================================================
# P3 - SLIDE 15 - Cas d'usage GenAI
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 3  |  L'IA Pero", "Cas d'usage & architecture GenAI", P3)
badge(s, "C5.1", "Cas d'usage GenAI", 0.52, 1.3, fill=EFREI_NAVY)
bullets(s, [
    ("Besoin  ",              "recommander un cocktail depuis une envie en langage naturel."),
    ("Pourquoi GenAI  ",      "les filtres echouent sur « frais, fruite, touche tropicale »."),
    ("Comprendre + produire  ","semantique (retrieval) ET generation (recette personnalisee)."),
    ("Conformite  ",          "thematique alternative validee par l'Annexe I du sujet."),
], 0.52, 1.9, 6.12, 3.22, size=14, gap=9, bullet_color=EFREI_PINK)
textbox(s, "Pipeline", 0.52, 5.2, 6, 0.32,
        size=14, color=P3, bold=True, font=FONT_H)
textbox(s, "Questionnaire -> SBERT -> cosinus -> guardrail 0.35 -> cache -> Gemini (RAG)",
        0.52, 5.52, 6.2, 0.78, size=13, color=TEXT_DARK, italic=True)
image_fit(s, SCREENS / "iapero-speakeasy-live.png", 6.88, 1.9, 5.9, 4.38, border=True)
caption(s, "Interface Speakeasy de L'IA Pero", 6.88, 6.32, 5.9)
footer(s, 15)

# ============================================================
# P3 - SLIDE 16 - Solution + DEMO
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 3  |  L'IA Pero", "Solution & demonstration", P3)
badge(s, "C5.2", "Solution GenAI", 0.52, 1.3, fill=EFREI_NAVY)
sp_d3 = rect(s, 9.52, 1.26, 3.28, 0.48, fill=EFREI_PINK,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, round_=0.3)
shape_text(sp_d3, "DEMO LIVE  -  90 s", size=12.5, color=WHITE)
bullets(s, [
    ("SBERT all-MiniLM-L6-v2  ",
     "local, gratuit, 384 dims, suffisant pour des phrases courtes."),
    ("Gemini free-tier  ",
     "generation uniquement, bornee par le contexte RAG."),
    ("Cache SQLite MD5  ",
     "1 appel API par generation -> conforme a la contrainte free-tier."),
    ("Guardrail  ",
     "similarite < 0.35 -> refus automatique des requetes hors-domaine."),
], 0.52, 1.92, 6.02, 3.3, size=13.5, gap=8, bullet_color=EFREI_PINK)
textbox(s, "Scenario : « quelque chose de frais et fruite » -> Top-5 + radar ; "
           "« repare mon velo » -> refus guardrail.",
        0.52, 5.3, 6.12, 0.88, size=12, color=TEXT_MUTED, italic=True)
image_fit(s, SCREENS / "iapero-recommendation-full.png", 6.88, 1.9, 5.9, 4.05, border=True)
button(s, "Ouvrir le repo ->", REPOS["iapero"], 6.88, 6.15, w=5.9, h=0.5, fill=P3)
footer(s, 16)

# ============================================================
# P3 - SLIDE 17 - Evaluation & risques
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Projet 3  |  L'IA Pero", "Evaluation & risques", P3)
badge(s, "C5.3", "Evaluation qualite", 0.52, 1.3, fill=EFREI_NAVY)
bullets(s, [
    ("Guardrail teste  ",
     "seuil 0.35 sur jeux de requetes hors-domaine documentes."),
    ("Tableau  ",
     "requete -> attendu -> score de similarite -> decision."),
    ("Ajustements  ",
     "seuil, Top-N, temperature Gemini, structure du prompt RAG."),
    ("Objectif tenu  ",
     "reponse < 3 s avec cache, < 8 s premier appel Gemini."),
], 0.52, 1.9, 6.12, 3.22, size=14, gap=9, bullet_color=EFREI_PINK)
rect(s, 0.52, 5.18, 6.12, 1.55,
     fill=RGBColor(0xFF, 0xF0, 0xF8), line=EFREI_PINK)
textbox(s, "Risques", 0.76, 5.3, 5, 0.32,
        size=13.5, color=EFREI_PINK, bold=True, font=FONT_H)
textbox(s, "Hallucination (mitigee RAG+cache)  |  dependance API (fallback+cache)  |  "
           "biais dataset  |  usage responsable de l'alcool.",
        0.76, 5.65, 5.7, 0.98, size=12, color=TEXT_DARK, line_spacing=1.12)
image_fit(s, SCREENS / "iapero-similarity-live.png", 6.88, 1.9, 5.9, 4.52, border=True)
caption(s, "Scores de similarite + decision du guardrail", 6.88, 6.48, 5.9)
footer(s, 17)

# ============================================================
# SLIDE 18 - Synthese competences
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Conclusion", "Synthese : competences -> preuves")
comps = [
    ("C1.1","Schema etoile PG"),    ("C1.2","Cassandra query-first"),
    ("C1.3","Medaillon Parquet"),   ("C1.4","Resilience testee"),
    ("C2.1","API JWT + quotas"),    ("C2.2","Kafka + micro-batch"),
    ("C2.3","Polars multi-sources"),("C2.4","Metriques pipeline"),
    ("C3.1","Prep fit-train"),      ("C3.2","Dashboard Streamlit"),
    ("C3.3","EDA correlations"),    ("C4.1","Strategie predictive"),
    ("C4.2","4 modeles testes"),    ("C4.3","Eval + CO2"),
    ("C5.1","Cas d'usage NL"),      ("C5.2","SBERT + RAG Gemini"),
    ("C5.3","Guardrail 0.35"),
]
cols = 4; cw_c = 3.05; chh = 0.92; gx = 0.12; gy = 0.14
x0, y0 = 0.52, 1.48
for i, (code, proof) in enumerate(comps):
    r_i = i // cols; c_i = i % cols
    x = x0 + c_i * (cw_c + gx); y = y0 + r_i * (chh + gy)
    sp = rect(s, x, y, cw_c, chh, fill=OFF_WHITE, line=RULE)
    rect(s, x, y, 0.09, chh, fill=EFREI_PINK)
    textbox(s, code,  x + 0.2, y + 0.1,  cw_c - 0.3, 0.34,
            size=14.5, color=EFREI_NAVY, bold=True, font=FONT_H)
    textbox(s, proof, x + 0.2, y + 0.46, cw_c - 0.3, 0.38,
            size=11.5, color=TEXT_DARK)
textbox(s, "17 competences C1.1 -> C5.3 couvertes, chacune avec une preuve dans le code.",
        0.52, 6.92, 12, 0.3, size=11.5, color=SUCCESS, bold=True)
footer(s, 18)

# ============================================================
# SLIDE 19 - Limites & ameliorations
# ============================================================
s = slide(); bg(s, WHITE)
header(s, "Conclusion", "Limites & ameliorations transverses")
items = [
    ("Haute disponibilite reelle",
     "Passer du mono-noeud a un cluster 3 noeuds (PG + Cassandra)."),
    ("Dimension temporelle",
     "LSTM sur sequences capteurs pour la maintenance predictive."),
    ("Base vectorielle",
     "pgvector / FAISS persistant au lieu de la matrice memoire pour L'IA Pero."),
    ("Monitoring & derive",
     "Detection de drift + reentrainement automatise en production."),
    ("CI/CD complet",
     "Pipeline GitHub Actions (tests + build + deploy) sur les 3 repos."),
]
y = 1.52
for titre, corps in items:
    rect(s, 0.52, y, 12.28, 0.95, fill=OFF_WHITE, line=RULE)
    rect(s, 0.52, y, 0.09, 0.95, fill=EFREI_PINK)
    textbox(s, titre, 0.8, y + 0.1,  4.2, 0.72, size=14.5, color=EFREI_NAVY,
            bold=True, font=FONT_H, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, corps, 5.18, y + 0.1, 7.4, 0.72, size=13, color=TEXT_DARK,
            anchor=MSO_ANCHOR.MIDDLE)
    y += 1.04
footer(s, 19)

# ============================================================
# SLIDE 20 - Conclusion (fond EFREI_DARK)
# ============================================================
s = slide(); bg(s, EFREI_DARK)
rect(s, 0, 0, 0.22, SH, fill=EFREI_PINK)
rect(s, 0, 0, SW, 0.08, fill=EFREI_PINK)
textbox(s, "Conclusion & apprentissages",
        0.55, 0.62, 11, 0.72, size=29, color=WHITE, bold=True, font=FONT_H)
rect(s, 0.55, 1.5, 5.5, 0.05, fill=EFREI_PINK)
bullets(s, [
    "Brancher tests et metriques des le jour 1, pas en consolidation",
    "Justifier chaque choix technique par un besoin metier concret",
    "Documenter les limites = posture professionnelle, pas un aveu",
    "Une chaine complete : data engineering -> data science -> GenAI",
], 0.55, 1.78, 11.5, 2.65, size=17, color=WHITE, gap=12,
    bullet_color=EFREI_PINK)
textbox(s, "Les 3 projets en un clic",
        0.55, 4.68, 11, 0.36, size=13.5,
        color=RGBColor(0x6B, 0xA3, 0xD8), bold=True, font=FONT_H)
button(s, "Urban Data Explorer ->",    REPOS["urban"],  0.55, 5.14, w=3.88, h=0.68,
       fill=WHITE, fg=EFREI_NAVY)
button(s, "Maintenance Predictive ->", REPOS["maint"],  4.58, 5.14, w=3.88, h=0.68,
       fill=WHITE, fg=P2)
button(s, "L'IA Pero ->",              REPOS["iapero"], 8.62, 5.14, w=3.88, h=0.68,
       fill=EFREI_PINK, fg=WHITE)
textbox(s, "Merci de votre attention  |  Adam Beloucif  &  Emilien Morice",
        0.55, 6.28, 11.5, 0.48, size=17.5, color=WHITE, bold=True, font=FONT_H)

# ============================================================
# ANNEXES backup (hors decompte)
# ============================================================
def annexe(title, imgs, key, col):
    s = slide(); bg(s, WHITE)
    header(s, "Annexe  |  Backup demo (hors decompte)", title, col)
    n = len(imgs)
    if n == 1:
        image_fit(s, imgs[0], 1.5, 1.42, 10.3, 5.28, border=True)
    else:
        w = (12.28 - 0.38 * (n - 1)) / n
        x = 0.52
        for img in imgs:
            image_fit(s, img, x, 1.52, w, 4.88, border=True)
            x += w + 0.38
    button(s, "Ouvrir le repo ->", REPOS[key], 5.38, 6.72, w=2.65, h=0.46, fill=col)


annexe("Urban Data Explorer  -  captures",
       [SCREENS / "urban-dashboard-live.png",
        SCREENS / "urban-district-selected.png",
        SCREENS / "ude-dark.png"],
       "urban", P1)
annexe("Maintenance Predictive  -  captures",
       [SCREENS / "maintenance-dashboard-live.png",
        SCREENS / "maintenance-diagnostic-live.png",
        SCREENS / "monthly_trends.png"],
       "maint", P2)
annexe("L'IA Pero  -  captures",
       [SCREENS / "iapero-speakeasy-live.png",
        SCREENS / "iapero-recommendation-full.png"],
       "iapero", P3)

prs.save(str(OUT))
print("OK ->", OUT)
print("slides:", len(prs.slides._sldIdLst))
